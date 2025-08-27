import io
from pptx import Presentation
from pptx.util import Inches

def build_ppt(slides, image_cycle=None):
    """
    Build PowerPoint presentation from slide dicts.

    slides: list of dicts with keys:
        - "title": str
        - "bullets": list[str]
        - "notes": str (optional)
        - "layout": str (optional)

    image_cycle: optional list of images (bytes) to reuse decoratively.
    """
    prs = Presentation()

    for idx, s in enumerate(slides):
        title = s.get("title", "")
        bullets = s.get("bullets", []) or []
        notes = s.get("notes", None)
        layout_hint = s.get("layout", "auto")

        # Decide layout
        need_body = len(bullets) > 0
        layout = _find_best_layout(prs, need_title=True, need_body=need_body)
        slide = prs.slides.add_slide(layout)

        # Title
        try:
            if slide.shapes.title:
                slide.shapes.title.text = title
        except Exception:
            pass

        # Body placeholder (bullets)
        try:
            body = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # TITLE
                    continue
                if shape.has_text_frame:
                    body = shape
                    break
            if body is not None and len(bullets) > 0:
                add_bullets(body.text_frame, bullets)
        except Exception:
            pass

        # Reuse a decorative image from template on content slides
        try:
            if image_cycle:
                blob = image_cycle[idx % len(image_cycle)]
                # Place near bottom-right with a modest size
                left = prs.slide_width - Inches(3)
                top = prs.slide_height - Inches(2)
                slide.shapes.add_picture(io.BytesIO(blob), left, top, width=Inches(2.5))
        except Exception:
            pass

        # Speaker notes
        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                # Some layouts may not have a notes slide yet
                try:
                    _ = slide.notes_slide
                    slide.notes_slide.notes_text_frame.text = notes
                except Exception:
                    pass

    return prs


def _find_best_layout(prs, need_title=True, need_body=True):
    """Pick the most appropriate layout available in the presentation template."""
    for layout in prs.slide_layouts:
        has_title = any(shape.placeholder_format.type == 1 for shape in layout.placeholders)
        has_body = any(shape.placeholder_format.type == 2 for shape in layout.placeholders)

        if need_title == has_title and need_body == has_body:
            return layout

    # fallback
    return prs.slide_layouts[0]


def add_bullets(text_frame, bullets):
    """Add bullet points to a slide body text frame."""
    text_frame.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
